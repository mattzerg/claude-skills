#!/usr/bin/env python3
"""consultant-deck — shadow-outline-first action-title deck producer."""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path.home() / ".claude/skills/_consultant/python"))
from _bootstrap_helper import ensure_venv  # noqa: E402

ensure_venv(__file__)

import argparse  # noqa: E402
import re  # noqa: E402


def _extract_section(body: str, name: str) -> str:
    current = None
    buf: list[str] = []
    for line in body.splitlines():
        m = re.match(r"^##\s+(.+?)\s*$", line)
        if m:
            if current and current.lower() == name.lower():
                return "\n".join(buf).strip()
            current = m.group(1).strip()
            buf = []
        elif current:
            buf.append(line)
    if current and current.lower() == name.lower():
        return "\n".join(buf).strip()
    return ""


KEY_PAT = re.compile(r"^###\s+Key\s+(\d+)\s+[—-]\s+(.+?)$", re.MULTILINE)
SUPPORT_PAT = re.compile(r"^- \*\*S(\d+(?:\.\d+)?)\*\*\s+[—-]\s+(.+?)$", re.MULTILINE)
UPSTREAM_CITE_PAT = re.compile(r"`upstream:\s*([^`\n]+?)`")


def _resolve_support_content(claim: str, engagement_root: Path) -> dict | None:
    """Parse an `upstream: <path>` cite from a Minto support claim and resolve to
    a slide spec. Returns {"kind": "chart"|"table", ...} or None if no upstream
    or unresolvable."""
    from consultant_kit import frontmatter  # type: ignore

    m = UPSTREAM_CITE_PAT.search(claim)
    if not m:
        return None
    cite_raw = m.group(1).strip()
    # Strip backticks, surrounding quotes
    cite_raw = cite_raw.strip("`'\" ")

    # Try as relative path first
    candidates = [
        engagement_root / cite_raw,
        Path(cite_raw),
    ]
    target: Path | None = None
    for c in candidates:
        if c.exists():
            target = c
            break
    if target is None:
        return None

    if target.suffix.lower() in (".png", ".jpg", ".jpeg", ".svg"):
        # Use as chart directly (PNG preferred)
        png = target.with_suffix(".png") if target.suffix != ".png" else target
        if png.exists():
            return {"kind": "chart", "path": str(png)}
        return {"kind": "chart", "path": str(target)}

    if target.suffix.lower() == ".md":
        # Check for a paired PNG in 05-analysis/charts/ by slug-match
        try:
            fm, body = frontmatter.parse(target)
        except Exception:  # noqa: BLE001
            return None
        # Frontmatter may declare a chart path
        for key in ("chart_path", "chart"):
            cp = fm.get(key)
            if cp and Path(cp).exists():
                return {"kind": "chart", "path": str(cp)}
        # Look for any PNG in the same directory or in 05-analysis/charts/ that
        # shares a slug stem
        stem = target.stem
        for search_dir in (target.parent, engagement_root / "05-analysis/charts"):
            if search_dir.exists():
                for png in search_dir.glob(f"{stem}*.png"):
                    return {"kind": "chart", "path": str(png)}
        # Fallback: extract markdown table from body
        table = _extract_markdown_table(body)
        if table:
            return {"kind": "table", "path": str(target), "table_md": table}

    return None


def _scan_charts(engagement_root: Path) -> list[Path]:
    """Return all rendered chart PNGs in the engagement's analysis + frameworks dirs, sorted."""
    out: list[Path] = []
    for sub in ("05-analysis/charts", "frameworks"):
        d = engagement_root / sub
        if d.exists():
            out.extend(sorted(d.glob("*.png")))
    return out


def _extract_markdown_table(body: str) -> str | None:
    """Pull the largest markdown table from a body. Returns the raw table block
    (header + separator + rows), or None if no table found."""
    lines = body.splitlines()
    best_block: list[str] = []
    current: list[str] = []
    in_table = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            current.append(stripped)
            in_table = True
        else:
            if in_table:
                if len(current) > len(best_block):
                    best_block = current
                current = []
                in_table = False
    if in_table and len(current) > len(best_block):
        best_block = current
    if len(best_block) < 3:  # header + separator + ≥1 row
        return None
    # Drop any trailing rows that are all empty/dashes (clean up extra separators)
    while best_block and re.match(r"^[\|\s\-:]+$", best_block[-1] or ""):
        if len(best_block) > 2:  # keep at least header+sep
            best_block.pop()
        else:
            break
    return "\n".join(best_block)


def _scan_support_content(engagement_root: Path) -> list[dict]:
    """Return ordered list of support-slide content candidates:
    each item is either {"kind": "chart", "path": str} for a PNG, or
    {"kind": "table", "path": str, "table_md": str, "title_hint": str}
    for a framework markdown table (Pugh / RACI / BMC / SWOT etc.).

    Charts come first (PNGs from 05-analysis/charts + frameworks),
    then standalone framework markdown tables that don't have a paired PNG.
    """
    from consultant_kit import frontmatter  # type: ignore

    out: list[dict] = []
    charts = _scan_charts(engagement_root)
    chart_stems = {p.stem for p in charts}
    for p in charts:
        out.append({"kind": "chart", "path": str(p)})

    fw_dir = engagement_root / "frameworks"
    if fw_dir.exists():
        for md in sorted(fw_dir.glob("*.md")):
            if md.stem in chart_stems:
                # Already covered by the framework's PNG
                continue
            try:
                fm, body = frontmatter.parse(md)
            except Exception:  # noqa: BLE001
                continue
            table = _extract_markdown_table(body)
            if not table:
                continue
            title_hint = fm.get("framework") or md.stem
            out.append({"kind": "table", "path": str(md),
                        "table_md": table, "title_hint": title_hint})
    return out


def _scan_raid_risks(engagement_root: Path, *, top_n: int = 3) -> list[str]:
    """Pull HIGH-severity risk descriptions from the engagement's 07-raid.md."""
    raid_path = engagement_root / "07-raid.md"
    if not raid_path.exists():
        return []
    try:
        text = raid_path.read_text(encoding="utf-8")
    except Exception:  # noqa: BLE001
        return []
    risks: list[tuple[str, str]] = []
    in_risks = False
    for line in text.splitlines():
        if line.startswith("## Risks"):
            in_risks = True
            continue
        if line.startswith("## "):
            in_risks = False
            continue
        if not in_risks:
            continue
        if not line.startswith("|") or "---" in line or "| ID " in line:
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        if not cells or not cells[0] or cells[0] == "ID":
            continue
        rid, desc, sev = cells[0], (cells[1] if len(cells) > 1 else ""), (cells[2] if len(cells) > 2 else "")
        if rid.startswith("R") and desc and not desc.startswith("_["):
            risks.append((sev.lower(), desc))
    # Prefer HIGH, then MED, then anything
    risks.sort(key=lambda r: {"high": 0, "med": 1}.get(r[0], 2))
    return [d for _, d in risks[:top_n]]


def _scan_workplan_next_steps(engagement_root: Path, *, top_n: int = 3) -> list[str]:
    """Pull first N task names from the engagement's 07-workplan.md."""
    wp_path = engagement_root / "07-workplan.md"
    if not wp_path.exists():
        return []
    try:
        text = wp_path.read_text(encoding="utf-8")
    except Exception:  # noqa: BLE001
        return []
    steps: list[str] = []
    in_section = False
    for line in text.splitlines():
        if line.startswith("## Workplan"):
            in_section = True
            continue
        if line.startswith("## ") and in_section:
            in_section = False
            continue
        if not in_section or not line.startswith("|") or "---" in line or "| ID " in line:
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        if not cells or not cells[0] or cells[0] == "ID":
            continue
        tid, name = cells[0], (cells[1] if len(cells) > 1 else "")
        if tid.startswith("T") and name and not name.startswith("_["):
            owner = cells[2] if len(cells) > 2 else ""
            label = f"{name} ({owner})" if owner and owner != "—" else name
            steps.append(label)
    return steps[:top_n]


def _scan_source_citations(engagement_root: Path) -> list[dict]:
    """Walk every artifact in the engagement and collect:
    - frontmatter `source_citations` array
    - `[source: ...]` inline tags in body markdown
    Skips `[needs-verification]` and dedupes by (claim, source).
    """
    from consultant_kit import frontmatter  # type: ignore

    seen = set()
    out: list[dict] = []
    inline_src = re.compile(r"\[source:\s*([^\]]+)\]")

    for md in engagement_root.rglob("*.md"):
        try:
            fm, body = frontmatter.parse(md)
        except Exception:  # noqa: BLE001
            continue
        # Skip the storyline itself + minto (they reference these — don't double-count)
        if md.name in ("storyline.md", "06-synthesis-minto.md"):
            continue
        # Frontmatter source_citations
        for c in fm.get("source_citations") or []:
            src = c.get("source", "")
            if not src or "[needs-verification]" in src:
                continue
            key = (c.get("claim", "")[:80], src[:80])
            if key not in seen:
                seen.add(key)
                out.append(c)
        # Inline `[source: ...]` tags in body — each becomes a row with the line above as claim
        body_lines = body.splitlines()
        for i, line in enumerate(body_lines):
            for m in inline_src.finditer(line):
                src = m.group(1).strip()
                if not src:
                    continue
                # Use the line as the claim (cleaned of the tag)
                claim_line = inline_src.sub("", line).strip()
                claim_line = re.sub(r"^\s*[-*|]\s*", "", claim_line).strip()
                claim_line = re.sub(r"\*\*", "", claim_line)[:120]
                key = (claim_line[:80], src[:80])
                if key not in seen:
                    seen.add(key)
                    out.append({"claim": claim_line, "source": src, "url": ""})
    return out


def outline(args) -> int:
    from consultant_kit import frontmatter, io, cite  # type: ignore

    minto_path = Path(args.from_path)
    fm_up, body_up = frontmatter.parse(minto_path)
    engagement = args.engagement or fm_up.get("engagement")
    mode = args.mode or fm_up.get("mode") or "ops"

    governing = fm_up.get("governing") or "[draft governing thought]"
    gov_section = _extract_section(body_up, "governing thought")

    keys = KEY_PAT.findall(body_up)
    supports = SUPPORT_PAT.findall(body_up)

    # Group supports by key index (S1.x → key 1)
    by_key: dict[int, list[tuple[str, str]]] = {}
    for sid, sclaim in supports:
        kidx = int(sid.split(".")[0])
        by_key.setdefault(kidx, []).append((sid, sclaim))

    # Discover support content (charts + framework tables) + sources from the engagement
    engagement_root = io.engagement_dir(engagement, mode)
    support_pool = _scan_support_content(engagement_root)
    source_citations = _scan_source_citations(engagement_root)
    raid_risks = _scan_raid_risks(engagement_root)
    next_steps = _scan_workplan_next_steps(engagement_root)
    support_iter = iter(support_pool)

    def _next_support() -> dict | None:
        try:
            return next(support_iter)
        except StopIteration:
            return None

    slides = []
    # Title slide
    slides.append({
        "n": 1, "type": "title",
        "title": engagement.replace("-", " ").title(),
        "subtitle": "Consultant engagement readout",
        "eyebrow": "CONSULTANT ENGAGEMENT",
        "notes": f"Source storyline: {minto_path}",
    })
    # Executive summary
    slides.append({
        "n": 2, "type": "exec-summary",
        "title": "Executive summary",
        "governing": governing,
        "keys": [k for _, k in keys],
        "eyebrow": "EXECUTIVE SUMMARY",
        "notes": "Read top-down. Each key is a slide section below.",
    })
    # Key + support slides
    n = 3
    for kid, kclaim in keys:
        kidx = int(kid)
        # Section divider per key
        slides.append({
            "n": n, "type": "section-divider",
            "title": kclaim, "key_idx": kidx,
            "eyebrow": f"KEY {kidx} / {len(keys)}",
            "notes": "",
        })
        n += 1
        # Support slides: prefer explicit `upstream:` cite from Minto row;
        # fall back to next item from support_pool iterator
        for sid, sclaim in by_key.get(kidx, []):
            slide = {
                "n": n,
                "title": sclaim, "support_id": f"S{sid}", "key_idx": kidx,
                "eyebrow": f"ANALYSIS — KEY {kidx}",
                "caption": "",
                "source": "",
                "notes": "",
            }
            # First try resolving the upstream cite in the claim text
            resolved = _resolve_support_content(sclaim, engagement_root)
            support = resolved or _next_support()
            if support and support.get("kind") == "chart":
                slide["type"] = "chart"
                slide["chart_path"] = support["path"]
            elif support and support.get("kind") == "table":
                slide["type"] = "table"
                slide["table_md"] = support["table_md"]
                slide["notes"] = f"Source artifact: {support.get('path','')}"
            else:
                slide["type"] = "support"
                slide["chart_path"] = None
            slides.append(slide)
            n += 1
    # Recommendation slide — auto-fill from Minto governing + RAID risks + workplan
    rec_action = governing if governing and not governing.startswith("[draft") else "[draft the single action this engagement recommends]"
    slides.append({
        "n": n, "type": "recommendation",
        "title": f"Recommendation: {rec_action[:80].rstrip('.')}." if rec_action and not rec_action.startswith("[draft") else "Recommendation",
        "action": rec_action,
        "risks": raid_risks or ["[risk 1 — populate 07-raid.md to auto-fill]",
                                "[risk 2]", "[risk 3]"],
        "next_steps": next_steps or ["[step 1 — populate 07-workplan.md to auto-fill]",
                                     "[step 2]", "[step 3]"],
        "eyebrow": "RECOMMENDATION",
        "notes": f"Auto-pulled {len(raid_risks)} risks from RAID, {len(next_steps)} steps from workplan.",
    })
    n += 1
    # Appendix (auto-populated with discovered source citations)
    slides.append({
        "n": n, "type": "appendix-sources",
        "title": "Sources",
        "sources": source_citations,
        "eyebrow": "APPENDIX",
        "notes": f"{len(source_citations)} citations collected from engagement artifacts.",
    })

    # Lint
    findings = []
    if len(slides) > 35:
        findings.append(("HIGH", f"Deck has {len(slides)} slides — hard cap is 35; trim the Minto pyramid"))
    elif len(slides) > 25:
        findings.append(("MED", f"Deck has {len(slides)} slides — over 25 means Minto wasn't tight; consider trimming supports"))
    for s in slides:
        if s["type"] in ("section-divider", "chart", "support"):
            words = len(s["title"].split())
            if words < 5:
                findings.append(("MED", f"Slide {s['n']} title is {words} words — likely a topic, not an action title: {s['title']!r}"))
        if s["type"] == "support":
            findings.append(("MED", f"Slide {s['n']} has no chart_path or table_md attached — body will be empty"))

    # Client-mode citation gate
    if mode == "client" and cite.needs_verification(body_up):
        findings.append(("HIGH", "Storyline carries `[needs-verification]` tags — client-mode render will refuse until resolved"))

    body_lines = [
        f"# Storyline — {engagement}",
        "",
        f"**Source:** `{minto_path}`",
        f"**Mode:** {mode}",
        f"**Slides:** {len(slides)}",
        "",
        "## Governing thought",
        "",
        f"> {governing}",
        "",
        "## Slide list (action-title format)",
        "",
        "| # | Type | Action title |",
        "|---|---|---|",
    ]
    for s in slides:
        body_lines.append(f"| {s['n']} | {s['type']} | {s['title']} |")

    if findings:
        body_lines.append("")
        body_lines.append("## Lint findings")
        for sev, msg in findings:
            body_lines.append(f"- **{sev}** — {msg}")

    body_lines.append("")
    body_lines.append("## Approval gate")
    body_lines.append("")
    body_lines.append("Storyline must be approved before `render`. Edit titles above to taste, then re-run `consultant-deck render` on this file.")

    fm = frontmatter.envelope(
        engagement=engagement,
        slug=f"{io.slugify(engagement)}-storyline",
        skill="consultant-deck",
        inputs=[str(minto_path)],
        upstream=[str(minto_path)],
        extra={"mode": mode, "slides": slides, "approved": False},
    )

    out_root = io.engagement_dir(engagement, mode) / "08-deliverable"
    out_root.mkdir(parents=True, exist_ok=True)
    out_path = out_root / "storyline.md"
    frontmatter.write_md(out_path, fm, "\n".join(body_lines))
    print(f"wrote {out_path}")
    print(f"\n{len(slides)} slides scaffolded.")
    if findings:
        print("\nLint findings:")
        for sev, msg in findings:
            print(f"  {sev}: {msg}")
    print("\nNEXT: edit titles + flip `approved: true` in frontmatter, then `consultant-deck render`")
    return 0


def render(args) -> int:
    import datetime as _dt
    from consultant_kit import frontmatter, brand, cite, layout as _layout  # type: ignore
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches

    storyline_path = Path(args.path)
    fm, body = frontmatter.parse(storyline_path)
    mode = fm.get("mode") or fm.get("extra", {}).get("mode", "ops")
    slides_spec = fm.get("slides") or fm.get("extra", {}).get("slides") or []
    engagement = fm.get("engagement")

    if mode == "client" and cite.needs_verification(body):
        print("REFUSED: client-mode render blocked — storyline carries `[needs-verification]` tags. Resolve and re-approve.")
        return 2

    palette = brand.get(args.palette or "default")
    _layout.register_fonts()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(slides_spec)
    date = _dt.date.today().isoformat()
    ctx_base = {
        "engagement": engagement,
        "date": date,
        "slide_total": total,
    }

    for s in slides_spec:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        ctx = {**ctx_base, "slide_n": s.get("n") or (slides_spec.index(s) + 1)}
        _layout.dispatch(s.get("type", "support"), prs, slide, palette, s, ctx)

        # Speaker notes
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    # Output
    out_root = Path(storyline_path).parent
    deck_path = out_root / f"{engagement}-deck.pptx"
    prs.save(str(deck_path))
    print(f"wrote {deck_path} ({total} slides)")
    return 0


def main() -> int:
    p = argparse.ArgumentParser(prog="consultant-deck")
    sub = p.add_subparsers(dest="cmd", required=True)

    o = sub.add_parser("outline")
    o.add_argument("--from", dest="from_path", required=True)
    o.add_argument("--engagement", default=None)
    o.add_argument("--mode", choices=("client", "pm", "ops", "life"), default=None)
    o.set_defaults(func=outline)

    r = sub.add_parser("render")
    r.add_argument("path")
    r.add_argument("--target", choices=("pptx", "gslides", "gamma"), default="pptx")
    r.add_argument("--palette", choices=("default", "dark", "navy"), default="default")
    r.set_defaults(func=render)

    args = p.parse_args()
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
