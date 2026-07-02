"""pptx layout helpers for consultant-deck.

Centralizes slide chrome (eyebrow + action title + accent rule + footer band)
plus per-type body layouts. Every layout function takes
`(prs, slide, palette, spec)` and renders into an existing blank slide.

Layouts implemented:
- slide_title (slide 1)
- slide_exec_summary (slide 2)
- slide_section_divider
- slide_chart (image-embed body)
- slide_two_col
- slide_stats_strip
- slide_quote
- slide_table (markdown → pptx table)
- slide_recommendation
- slide_appendix_sources
- slide_support (fallback for legacy storylines; emits placeholder + lint warning)

Phase 1 ships chrome + slide_title + slide_exec_summary + slide_chart +
slide_support fallback. Phase 3 fills out the rest.
"""
from __future__ import annotations

import datetime as _dt
from pathlib import Path
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import brand

FONT_STACK = "Space Grotesk"
FONT_FALLBACK = "Helvetica Neue"
FONT_MONO = "Space Mono"

# Slide canvas: 13.333 x 7.5 in
CANVAS_W = Inches(13.333)
CANVAS_H = Inches(7.5)

# Chrome zones
EYEBROW_TOP = Inches(0.4)
EYEBROW_H = Inches(0.25)
TITLE_TOP = Inches(0.75)
TITLE_H = Inches(1.0)
ACCENT_RULE_TOP = Inches(1.85)
ACCENT_RULE_W = Inches(2.0)
ACCENT_RULE_H = Emu(36000)  # ~0.04 in
BODY_TOP = Inches(2.15)
BODY_H = Inches(4.85)
FOOTER_TOP = Inches(7.05)
FOOTER_H = Inches(0.35)

# Side margins
LMARGIN = Inches(0.6)
RMARGIN = Inches(0.6)
INNER_W = Inches(12.13)  # CANVAS_W - 2*LMARGIN


def _hex(s: str) -> RGBColor:
    return RGBColor.from_string(s.lstrip("#"))


def _set_run(run, *, text: str, font: str = FONT_STACK, size: int = 11,
             bold: bool = False, italic: bool = False, color: str | None = None,
             color_rgb: RGBColor | None = None, all_caps: bool = False) -> None:
    """Apply standard run attributes. Uses run.text rather than appending."""
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color_rgb is not None:
        run.font.color.rgb = color_rgb
    elif color is not None:
        run.font.color.rgb = _hex(color)
    # python-pptx doesn't expose all-caps directly; we uppercase manually
    if all_caps:
        run.text = text.upper()


def paint_background(slide, palette: brand.Palette) -> None:
    """Set the slide background to the palette's paper color."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex(palette.paper)


# --- Chrome -----------------------------------------------------------------

def draw_eyebrow(slide, palette: brand.Palette, label: str) -> None:
    """Tiny ALL-CAPS tracked label in accent color above the action title."""
    if not label:
        return
    tx = slide.shapes.add_textbox(LMARGIN, EYEBROW_TOP, INNER_W, EYEBROW_H)
    tf = tx.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(
        run, text=label, font=FONT_MONO, size=8, bold=True,
        color=palette.accent_primary_alt if palette.name == "zerg-default" else palette.accent_primary,
        all_caps=True,
    )


def draw_action_title(slide, palette: brand.Palette, title: str, *, size: int = 22) -> None:
    """Action title — a complete claim, not a topic."""
    tx = slide.shapes.add_textbox(LMARGIN, TITLE_TOP, INNER_W, TITLE_H)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(run, text=title, size=size, bold=True, color=palette.text)


def draw_accent_rule(slide, palette: brand.Palette, *,
                     top=ACCENT_RULE_TOP, width=ACCENT_RULE_W) -> None:
    """Short accent rule under the title."""
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LMARGIN, top, width, ACCENT_RULE_H)
    rule.fill.solid()
    rule.fill.fore_color.rgb = _hex(palette.accent_primary)
    rule.line.fill.background()


def draw_footer(slide, palette: brand.Palette, *,
                engagement: str, date: str, slide_n: int, slide_total: int) -> None:
    """Footer band: engagement (left) | date (center) | n/total (right)."""
    # Accent thin rule above
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LMARGIN, FOOTER_TOP - Emu(80000), INNER_W, Emu(12000),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = _hex(palette.rule_gray)
    rule.line.fill.background()

    third = Inches(4.04)
    cells = [
        (LMARGIN, engagement.replace("-", " "), PP_ALIGN.LEFT),
        (LMARGIN + third, date, PP_ALIGN.CENTER),
        (LMARGIN + third * 2, f"{slide_n} / {slide_total}", PP_ALIGN.RIGHT),
    ]
    for left, text, align in cells:
        tx = slide.shapes.add_textbox(left, FOOTER_TOP, third, FOOTER_H)
        tf = tx.text_frame
        tf.margin_left = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        _set_run(run, text=text, size=8, color=palette.mid_gray)


def chrome(slide, palette: brand.Palette, *,
           engagement: str, date: str, slide_n: int, slide_total: int,
           eyebrow: str, title: str, title_size: int = 22) -> None:
    """Apply standard chrome to a slide: eyebrow + title + accent rule + footer."""
    paint_background(slide, palette)
    draw_eyebrow(slide, palette, eyebrow)
    draw_action_title(slide, palette, title, size=title_size)
    draw_accent_rule(slide, palette)
    draw_footer(slide, palette,
                engagement=engagement, date=date,
                slide_n=slide_n, slide_total=slide_total)


# --- Per-type layouts -------------------------------------------------------

def slide_title(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Slide 1. Engagement name big; subtitle + date below."""
    paint_background(slide, palette)
    # Big eyebrow
    tx = slide.shapes.add_textbox(LMARGIN, Inches(2.0), INNER_W, Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    _set_run(
        r, text="CONSULTANT ENGAGEMENT", font=FONT_MONO, size=12,
        bold=True, color=palette.accent_primary, all_caps=True,
    )
    # Title (engagement name)
    tx = slide.shapes.add_textbox(LMARGIN, Inches(2.6), INNER_W, Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    _set_run(r, text=spec.get("title", ""), size=44, bold=True, color=palette.text)
    # Subtitle
    if spec.get("subtitle"):
        tx = slide.shapes.add_textbox(LMARGIN, Inches(4.7), INNER_W, Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, text=spec["subtitle"], size=16, color=palette.mid_gray)
    # Date (bottom-left)
    tx = slide.shapes.add_textbox(LMARGIN, Inches(6.6), INNER_W, Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, text=ctx.get("date", _dt.date.today().isoformat()), size=11, color=palette.mid_gray)
    # Accent rule
    draw_accent_rule(slide, palette, top=Inches(2.4), width=Inches(1.4))


def slide_exec_summary(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Governing thought + numbered keys."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow="EXECUTIVE SUMMARY", title=spec.get("title", "Executive summary"),
           title_size=24)

    # Governing thought
    gov = spec.get("governing", "")
    if gov:
        tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, INNER_W, Inches(1.0))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        _set_run(r, text=gov, size=16, color=palette.text, italic=True)

    # Numbered keys
    keys = spec.get("keys", []) or []
    if keys:
        tx = slide.shapes.add_textbox(LMARGIN, Inches(3.4), INNER_W, Inches(3.4))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, k in enumerate(keys):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(8)
            # Number in accent color
            num_run = p.add_run()
            _set_run(num_run, text=f"{i+1}.  ", size=15, bold=True, color=palette.accent_primary)
            txt_run = p.add_run()
            _set_run(txt_run, text=k, size=15, color=palette.text)


def slide_chart(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Embed a chart PNG with caption + source line."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow=spec.get("eyebrow", "ANALYSIS"),
           title=spec.get("title", ""))

    chart_path = spec.get("chart_path")
    if not chart_path or not Path(chart_path).exists():
        # Fall back to a placeholder with explicit warning
        tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, INNER_W, Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(
            r, text=f"[chart missing: {chart_path or 'no chart_path set on slide spec'}]",
            size=11, color=palette.mid_gray, italic=True,
        )
        return

    # Image: top 4.4in of the body region, centered horizontally
    img_top = BODY_TOP
    img_h = Inches(4.4)
    img_w = Inches(10.0)
    img_left = (CANVAS_W - img_w) / 2
    slide.shapes.add_picture(str(chart_path), img_left, img_top,
                             width=img_w, height=img_h)

    # Caption (≤2 lines, just below chart)
    caption = spec.get("caption", "")
    if caption:
        tx = slide.shapes.add_textbox(LMARGIN, img_top + img_h + Inches(0.05),
                                      INNER_W, Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        _set_run(r, text=caption, size=10, color=palette.mid_gray, italic=True)

    # Source line (right-aligned, just above footer)
    source = spec.get("source", "")
    if source:
        tx = slide.shapes.add_textbox(LMARGIN, Inches(6.6), INNER_W, Inches(0.3))
        p = tx.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        _set_run(r, text=f"Source: {source}", size=8, color=palette.mid_gray)


def slide_section_divider(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Big eyebrow + chapter title; used at phase transitions."""
    paint_background(slide, palette)
    tx = slide.shapes.add_textbox(LMARGIN, Inches(3.0), INNER_W, Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(
        r, text=spec.get("eyebrow", "SECTION"), font=FONT_MONO, size=14,
        bold=True, color=palette.accent_primary, all_caps=True,
    )
    tx = slide.shapes.add_textbox(LMARGIN, Inches(3.7), INNER_W, Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, text=spec.get("title", ""), size=40, bold=True, color=palette.text)
    draw_accent_rule(slide, palette, top=Inches(3.5), width=Inches(1.4))
    draw_footer(slide, palette,
                engagement=ctx["engagement"], date=ctx["date"],
                slide_n=ctx["slide_n"], slide_total=ctx["slide_total"])


def slide_two_col(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """1.6fr (text) | 1fr (chart or text). Phase 3 fills out fully."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow=spec.get("eyebrow", "ANALYSIS"),
           title=spec.get("title", ""))

    total_w = INNER_W
    left_w = Inches(7.46)  # ~1.6/2.6 × 12.13
    right_w = Inches(4.0)  # ~1/2.6 × 12.13
    gap = Inches(0.27)

    # Left column: bullets/prose
    bullets = spec.get("left_bullets") or []
    if bullets:
        tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, left_w, BODY_H)
        tf = tx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            r = p.add_run()
            _set_run(r, text=f"•  {b}", size=12, color=palette.text)

    # Right column: chart_path or text
    if spec.get("right_chart_path") and Path(spec["right_chart_path"]).exists():
        slide.shapes.add_picture(
            str(spec["right_chart_path"]),
            LMARGIN + left_w + gap, BODY_TOP,
            width=right_w, height=Inches(3.6),
        )
    elif spec.get("right_text"):
        tx = slide.shapes.add_textbox(LMARGIN + left_w + gap, BODY_TOP, right_w, BODY_H)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        _set_run(r, text=spec["right_text"], size=12, color=palette.text)


def slide_stats_strip(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """3–4 equal-width cells with stat + unit + caption."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow=spec.get("eyebrow", "KEY METRICS"),
           title=spec.get("title", ""))

    stats = spec.get("stats", []) or []
    if not stats:
        return
    n = min(len(stats), 4)
    gap = Inches(0.3)
    cell_w = Inches((12.13 - 0.3 * (n - 1)) / n)
    cell_top = Inches(3.0)
    cell_h = Inches(3.0)
    for i, s in enumerate(stats[:n]):
        left = LMARGIN + i * (cell_w + gap)
        # Stat number (big)
        tx = slide.shapes.add_textbox(left, cell_top, cell_w, Inches(1.2))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        _set_run(r, text=s.get("value", "—"), size=44, bold=True, color=palette.accent_primary)
        # Unit (eyebrow)
        if s.get("unit"):
            tx = slide.shapes.add_textbox(left, cell_top + Inches(1.3), cell_w, Inches(0.3))
            p = tx.text_frame.paragraphs[0]
            r = p.add_run()
            _set_run(
                r, text=s["unit"], font=FONT_MONO, size=9, bold=True,
                color=palette.accent_primary_alt if palette.name == "zerg-default" else palette.accent_primary,
                all_caps=True,
            )
        # Caption
        if s.get("caption"):
            tx = slide.shapes.add_textbox(left, cell_top + Inches(1.7), cell_w, Inches(1.2))
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            _set_run(r, text=s["caption"], size=11, color=palette.text)


def slide_quote(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Pull quote (large italic) + attribution."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow=spec.get("eyebrow", "VOICE OF CUSTOMER"),
           title=spec.get("title", ""))

    quote = spec.get("quote", "")
    attribution = spec.get("attribution", "")
    if quote:
        # Big accent quote mark
        tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, Inches(0.8), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, text="“", size=72, bold=True, color=palette.accent_primary)
        # Quote
        tx = slide.shapes.add_textbox(LMARGIN + Inches(0.7), BODY_TOP + Inches(0.2),
                                      INNER_W - Inches(0.7), Inches(3.4))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        _set_run(r, text=quote, size=24, italic=True, color=palette.text)
    if attribution:
        tx = slide.shapes.add_textbox(LMARGIN + Inches(0.7), Inches(6.2),
                                      INNER_W - Inches(0.7), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, text=f"— {attribution}", size=12, color=palette.mid_gray)


def md_table_to_rows(md: str) -> list[list[str]]:
    """Parse a markdown table into a 2D list of cells."""
    rows = []
    for line in md.splitlines():
        line = line.strip()
        if not line.startswith("|"):
            continue
        if line.startswith("|---") or set(line.replace("|", "").strip()) <= set("-: "):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        if cells:
            rows.append(cells)
    return rows


def slide_table(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Markdown table → PPTX table with banded rows + accent header."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow=spec.get("eyebrow", "DATA"),
           title=spec.get("title", ""))

    rows = md_table_to_rows(spec.get("table_md", ""))
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    # Normalize
    for r in rows:
        while len(r) < n_cols:
            r.append("")

    width = Inches(12.0)
    height = min(Inches(0.45) * n_rows, Inches(4.5))
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         (CANVAS_W - width) / 2, BODY_TOP,
                                         width, height)
    table = table_shape.table
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            if ri == 0:
                _set_run(r, text=cell_text, size=10, bold=True, color=palette.paper)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex(palette.accent_primary)
            else:
                _set_run(r, text=cell_text, size=10, color=palette.text)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex(palette.paper if ri % 2 else palette.rule_gray)


def slide_recommendation(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Action + risks + next steps."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow="RECOMMENDATION",
           title=spec.get("title", ""))

    # Action box (large accent)
    action = spec.get("action", "")
    if action:
        tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, INNER_W, Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        _set_run(r, text=action, size=20, bold=True, color=palette.accent_primary)

    # Two columns: risks (left), next steps (right)
    col_w = Inches(5.9)
    gap = Inches(0.33)
    col_top = Inches(3.7)

    risks = spec.get("risks", []) or []
    if risks:
        tx = slide.shapes.add_textbox(LMARGIN, col_top - Inches(0.3), col_w, Inches(0.3))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(
            r, text="RISKS", font=FONT_MONO, size=9, bold=True,
            color=palette.accent_primary_alt if palette.name == "zerg-default" else palette.accent_primary,
            all_caps=True,
        )
        tx = slide.shapes.add_textbox(LMARGIN, col_top, col_w, Inches(3.0))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, risk in enumerate(risks[:3]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            r = p.add_run()
            _set_run(r, text=f"•  {risk}", size=12, color=palette.text)

    next_steps = spec.get("next_steps", []) or []
    if next_steps:
        tx = slide.shapes.add_textbox(LMARGIN + col_w + gap, col_top - Inches(0.3),
                                      col_w, Inches(0.3))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(
            r, text="NEXT STEPS", font=FONT_MONO, size=9, bold=True,
            color=palette.accent_primary_alt if palette.name == "zerg-default" else palette.accent_primary,
            all_caps=True,
        )
        tx = slide.shapes.add_textbox(LMARGIN + col_w + gap, col_top, col_w, Inches(3.0))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, step in enumerate(next_steps[:3]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            r = p.add_run()
            _set_run(r, text=f"{i+1}.  {step}", size=12, color=palette.text)


def slide_appendix_sources(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Auto-generated sources list from upstream source_citations."""
    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow="APPENDIX",
           title=spec.get("title", "Sources"))

    sources = spec.get("sources", []) or []
    if not sources:
        tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, INNER_W, Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, text="No source citations collected upstream.", size=11,
                 color=palette.mid_gray, italic=True)
        return

    tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, INNER_W, BODY_H)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, s in enumerate(sources[:20]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        claim = s.get("claim", "")
        source = s.get("source", "")
        url = s.get("url", "")
        text = f"[{i+1}] {claim} — {source}"
        if url:
            text += f" ({url})"
        r = p.add_run()
        _set_run(r, text=text, size=10, color=palette.text)


def slide_support(prs, slide, palette: brand.Palette, spec: dict, ctx: dict) -> None:
    """Fallback for legacy slides without explicit chart/table content.
    Emits a clear placeholder that will be visible (and flagged by the lint)."""
    # If chart_path is set, route to slide_chart
    if spec.get("chart_path"):
        return slide_chart(prs, slide, palette, spec, ctx)
    if spec.get("table_md"):
        return slide_table(prs, slide, palette, spec, ctx)

    chrome(slide, palette,
           engagement=ctx["engagement"], date=ctx["date"],
           slide_n=ctx["slide_n"], slide_total=ctx["slide_total"],
           eyebrow=spec.get("eyebrow", "ANALYSIS"),
           title=spec.get("title", ""))

    # Visible "BODY MISSING" marker so reviewer can't miss it
    tx = slide.shapes.add_textbox(LMARGIN, BODY_TOP, INNER_W, Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, text="[body missing — attach chart_path or table_md to this slide spec]",
             size=12, color=palette.mid_gray, italic=True)


LAYOUTS = {
    "title": slide_title,
    "exec-summary": slide_exec_summary,
    "section-divider": slide_section_divider,
    "chart": slide_chart,
    "two-col": slide_two_col,
    "stats-strip": slide_stats_strip,
    "quote": slide_quote,
    "table": slide_table,
    "recommendation": slide_recommendation,
    "appendix": slide_appendix_sources,
    "appendix-sources": slide_appendix_sources,
    "support": slide_support,
    "key-section": slide_section_divider,
}


def dispatch(layout_name: str, prs, slide, palette: brand.Palette,
             spec: dict, ctx: dict) -> None:
    """Resolve and call the layout function for `layout_name`. Falls back to slide_support."""
    fn = LAYOUTS.get(layout_name, slide_support)
    fn(prs, slide, palette, spec, ctx)


def register_fonts() -> None:
    """Register Space Grotesk + Space Mono with matplotlib if available locally.
    Searches typical macOS install paths + ~/.claude/skills/_consultant/python/fonts/.
    Silent no-op if fonts not found (matplotlib falls back to sans-serif)."""
    try:
        import matplotlib.font_manager as fm  # type: ignore
    except ImportError:
        return
    candidates = [
        Path.home() / ".claude/skills/_consultant/python/fonts",
        Path.home() / "Library/Fonts",
        Path("/Library/Fonts"),
        Path("/System/Library/Fonts/Supplemental"),
    ]
    for d in candidates:
        if not d.exists():
            continue
        for f in d.glob("Space*Grotesk*.ttf"):
            try:
                fm.fontManager.addfont(str(f))
            except Exception:  # noqa: BLE001
                pass
        for f in d.glob("Space*Grotesk*.otf"):
            try:
                fm.fontManager.addfont(str(f))
            except Exception:  # noqa: BLE001
                pass
