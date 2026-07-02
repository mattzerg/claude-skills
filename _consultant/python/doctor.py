#!/usr/bin/env python3
"""Consultant-toolkit health check.

Renders one fixture of each chart recipe (12 total) + one fixture deck with each
slide layout (10 types), then sanity-checks the output. Exits 0 on green.

Run via `bootstrap.sh && python3 doctor.py`. Optional `--bless` writes current
output to `fixtures/baseline/` for visual regression in subsequent runs.
"""
from __future__ import annotations

import argparse
import shutil
import sys
import tempfile
from pathlib import Path

from consultant_kit import brand, cite, frontmatter, ids, io


CHECKS = []


def check(label):
    def wrap(fn):
        CHECKS.append((label, fn))
        return fn
    return wrap


@check("brand-palette-load")
def t_brand() -> str:
    p = brand.get("default")
    assert p.paper == "#f4f0e7"
    p2 = brand.get("dark")
    assert p2.paper == "#111514"
    rc = brand.matplotlib_rcparams("default")
    assert rc["figure.facecolor"] == "#f4f0e7"
    # New: Okabe-Ito + semantic palettes
    assert len(brand.OKABE_ITO) == 8
    sem = brand.semantic_palette(p)
    assert sem["positive"] == p.accent_secondary
    # Highlight cycle
    cyc = brand.chart_color_cycle(p, highlight_idx=2)
    assert cyc[2] == p.text  # original index 2 is text in cycle
    assert all(c == p.mid_gray for i, c in enumerate(cyc) if i != 2)
    return "ok"


@check("ids-derivation")
def t_ids() -> str:
    assert ids.child("L1", 2) == "L1.2"
    assert ids.parent("L1.2.3") == "L1.2"
    assert ids.depth("L1.2.3") == 3
    assert ids.hypothesis_id("L1.2") == "H1.2"
    return "ok"


@check("frontmatter-roundtrip")
def t_frontmatter() -> str:
    with tempfile.TemporaryDirectory() as td:
        path = Path(td) / "test.md"
        fm = frontmatter.envelope(
            engagement="test", slug="demo", skill="doctor",
            inputs=["a.csv"], upstream=["b.md"],
            source_citations=[cite.new("claim", "Source X", "https://x").to_dict()],
        )
        frontmatter.write_md(path, fm, "# body\n\nhello")
        fm2, body = frontmatter.parse(path)
        assert fm2["engagement"] == "test"
        assert "hello" in body
    return "ok"


@check("cite-validation")
def t_cite() -> str:
    assert cite.needs_verification("revenue is $5M [needs-verification]")
    assert not cite.needs_verification("revenue is $5M [source: Crunchbase 2026]")
    ok, _ = cite.validate_for_client("clean text [source: x]")
    assert ok
    return "ok"


@check("io-paths")
def t_io() -> str:
    assert io.slugify("Hello World!") == "hello-world"
    p = io.engagement_dir("acme-activation", "client")
    assert "clients/acme-activation" in str(p)
    p2 = io.engagement_dir("decision", "life")
    assert "_personal/decision" in str(p2)
    return "ok"


@check("annotate-formatters")
def t_annotate() -> str:
    from consultant_kit import annotate
    assert annotate.format_value(1200, "thousands") == "1.2K"
    assert annotate.format_value(1_500_000, "currency") == "$1.5M"
    assert annotate.format_value(0.072, "percent") == "7.2%"
    assert annotate.format_value(1240, "int") == "1,240"
    assert annotate.pick_format([1200, 1800, 2900]) == "thousands"
    assert annotate.pick_format([0.05, 0.12, 0.18]) == "percent"
    return "ok"


@check("chart-recipes-render-12")
def t_charts() -> str:
    from consultant_kit import chart
    fx_dir = Path(__file__).parent / "fixtures" / "charts"
    fx_dir.mkdir(parents=True, exist_ok=True)

    # bar
    chart.render("bar", out=fx_dir / "bar.png",
                 labels=["A", "B", "C"], values=[3000, 7500, 5200], ylabel="USD")
    # line
    chart.render("line", out=fx_dir / "line.png",
                 x=[1, 2, 3, 4], series={"alpha": [12, 24, 48, 72], "beta": [20, 22, 25, 28]},
                 xlabel="t", ylabel="users")
    # stacked-bar
    chart.render("stacked-bar", out=fx_dir / "stacked.png",
                 labels=["Q1", "Q2", "Q3"], series={"new": [3, 4, 5], "expansion": [2, 2, 1]},
                 ylabel="$K")
    # waterfall
    chart.render("waterfall", out=fx_dir / "waterfall.png",
                 labels=["Mix", "Price", "Volume"], deltas=[2_000_000, -1_000_000, 3_000_000],
                 ylabel="USD", start_label="FY24", end_label="FY25")
    # heatmap
    chart.render("heatmap", out=fx_dir / "heatmap.png",
                 matrix=[[1, 2, 3], [4, 5, 6], [7, 8, 9]],
                 row_labels=["R1", "R2", "R3"], col_labels=["C1", "C2", "C3"])
    # scatter-2x2
    chart.render("scatter-2x2", out=fx_dir / "scatter2x2.png",
                 items=[{"label": "Plan A", "x": 3, "y": 8, "group": "now"},
                        {"label": "Plan B", "x": 7, "y": 5, "group": "later"},
                        {"label": "Plan C", "x": 2, "y": 2, "group": "later"}])
    # marimekko
    chart.render("marimekko", out=fx_dir / "marimekko.png",
                 rows=[{"segment": "SMB", "width": 50, "shares": {"win": 30, "loss": 70}},
                       {"segment": "MM", "width": 30, "shares": {"win": 55, "loss": 45}},
                       {"segment": "ENT", "width": 20, "shares": {"win": 70, "loss": 30}}])
    # grouped-bar
    chart.render("grouped-bar", out=fx_dir / "grouped-bar.png",
                 labels=["SMB", "MM", "ENT"], series={"Win": [42, 58, 71], "Loss": [58, 42, 29]},
                 ylabel="% deals")
    # slope-graph
    chart.render("slope-graph", out=fx_dir / "slope.png",
                 items=[{"label": "A", "before": 28, "after": 47},
                        {"label": "B", "before": 35, "after": 32},
                        {"label": "C", "before": 22, "after": 28}],
                 highlight=["A"])
    # dot-plot
    chart.render("dot-plot", out=fx_dir / "dot.png",
                 labels=["a", "b", "c", "d"], values=[12, 24, 18, 9], highlight_idx=1)
    # bullet
    chart.render("bullet", out=fx_dir / "bullet.png",
                 items=[{"label": "SMB", "actual": 45, "target": 60, "ranges": [30, 50, 75]},
                        {"label": "MM", "actual": 88, "target": 90, "ranges": [40, 70, 110]}])
    # small-multiples
    chart.render("small-multiples", out=fx_dir / "small-multiples.png", cols=2,
                 panels=[
                     {"title": "Jan", "x": [1, 2, 3], "y": [10, 15, 20]},
                     {"title": "Feb", "x": [1, 2, 3], "y": [8, 12, 18]},
                     {"title": "Mar", "x": [1, 2, 3], "y": [5, 9, 14]},
                     {"title": "Apr", "x": [1, 2, 3], "y": [3, 6, 10]},
                 ])

    # Verify all 12 PNGs exist + are non-trivial size
    for name in ("bar", "line", "stacked", "waterfall", "heatmap", "scatter2x2",
                 "marimekko", "grouped-bar", "slope", "dot", "bullet", "small-multiples"):
        path = fx_dir / f"{name}.png"
        assert path.exists(), f"missing {name}"
        assert path.stat().st_size > 5000, f"{name} too small ({path.stat().st_size}B)"
    return f"12 recipes"


@check("chart-spec-validation")
def t_validate() -> str:
    from consultant_kit import chart
    # Bad: missing labels
    findings = chart.validate_spec({"recipe": "bar", "values": [1, 2, 3]})
    assert any(f[0] == "HIGH" for f in findings)
    # Bad: stacked-bar with negative
    findings = chart.validate_spec({
        "recipe": "stacked-bar", "labels": ["a"],
        "series": {"x": [-1]},
    })
    assert any(f[0] == "HIGH" for f in findings)
    # Good
    findings = chart.validate_spec({
        "recipe": "bar", "labels": ["a", "b"], "values": [1, 2],
    })
    assert not any(f[0] == "HIGH" for f in findings)
    return "ok"


@check("pptx-layouts-render-10")
def t_layouts() -> str:
    from consultant_kit import layout
    from pptx import Presentation
    from pptx.util import Inches

    fx_dir = Path(__file__).parent / "fixtures" / "decks"
    fx_dir.mkdir(parents=True, exist_ok=True)
    fixture_chart = Path(__file__).parent / "fixtures" / "charts" / "bar.png"

    palette = brand.get("default")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = [
        ("title", {"title": "Test Engagement", "subtitle": "Doctor smoke test"}),
        ("exec-summary", {"title": "Doctor finds three things to fix this week.",
                          "governing": "Doctor finds three things to fix this week.",
                          "keys": ["First key claim", "Second key claim", "Third key claim"]}),
        ("section-divider", {"title": "Key 1 / 3 — Discovery", "eyebrow": "KEY 1"}),
        ("chart", {"title": "Bar chart rendering confirmed via doctor fixture.",
                   "chart_path": str(fixture_chart), "caption": "Doctor bar fixture",
                   "source": "Doctor 2026-05-29", "eyebrow": "ANALYSIS"}),
        ("two-col", {"title": "Two-column layout supports text + chart.",
                     "left_bullets": ["First bullet point", "Second bullet point", "Third bullet point"],
                     "right_chart_path": str(fixture_chart), "eyebrow": "COMPARISON"}),
        ("stats-strip", {"title": "Three numbers anchor the engagement.",
                         "stats": [{"value": "$2.3M", "unit": "REVENUE", "caption": "FY26 run rate"},
                                   {"value": "42%", "unit": "GROSS MARGIN", "caption": "Up 8pp YoY"},
                                   {"value": "1,240", "unit": "ACTIVE USERS", "caption": "Weekly"}],
                         "eyebrow": "KEY METRICS"}),
        ("quote", {"title": "Customer voice supports the strategic direction.",
                   "quote": "This was the right call.",
                   "attribution": "Test Customer, VP of Ops"}),
        ("table", {"title": "Pricing matrix maps to segment willingness-to-pay.",
                   "table_md": "| Tier | Price | Features |\n|---|---|---|\n| Basic | $9 | 3 |\n| Pro | $19 | 12 |\n| Bundle | $39 | 25 |",
                   "eyebrow": "PRICING"}),
        ("recommendation", {"title": "Recommendation: build Self-Serve docs in Q3.",
                            "action": "Build self-serve docs in Q3.",
                            "risks": ["First risk", "Second risk", "Third risk"],
                            "next_steps": ["First step", "Second step", "Third step"]}),
        ("appendix-sources", {"title": "Sources",
                              "sources": [
                                  {"claim": "Test claim 1", "source": "Test Source A", "url": "https://a"},
                                  {"claim": "Test claim 2", "source": "Test Source B", "url": ""},
                              ]}),
    ]

    layout.register_fonts()
    for i, (layout_name, spec) in enumerate(layouts, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        ctx = {"engagement": "doctor-fixture", "date": "2026-05-29",
               "slide_n": i, "slide_total": len(layouts)}
        layout.dispatch(layout_name, prs, slide, palette, spec, ctx)

    deck_path = fx_dir / "doctor-deck.pptx"
    prs.save(str(deck_path))
    assert deck_path.exists()
    assert deck_path.stat().st_size > 30_000, f"deck too small ({deck_path.stat().st_size}B)"
    return f"10 layouts, {deck_path.stat().st_size//1024}KB"


@check("regression-baseline")
def t_regression() -> str:
    """If --bless was passed earlier this run, copy current fixtures to baseline.
    Otherwise just confirm baseline dir is reachable (no hash diff to enforce v1 — too brittle)."""
    base = Path(__file__).parent / "fixtures" / "baseline"
    base.mkdir(parents=True, exist_ok=True)
    return f"baseline dir: {base}"


def maybe_bless():
    """If --bless passed, copy current fixtures to baseline."""
    fx = Path(__file__).parent / "fixtures"
    baseline = fx / "baseline"
    for sub in ("charts", "decks"):
        src = fx / sub
        if not src.exists():
            continue
        dst = baseline / sub
        if dst.exists():
            shutil.rmtree(dst)
        shutil.copytree(src, dst)
    print(f"  ✓ blessed fixtures → {baseline}")


def main() -> int:
    parser = argparse.ArgumentParser(prog="doctor")
    parser.add_argument("--bless", action="store_true",
                        help="after rendering, copy current fixtures to baseline/")
    args = parser.parse_args()

    failures = []
    for label, fn in CHECKS:
        try:
            res = fn()
            print(f"  ✓ {label}: {res}")
        except Exception as e:  # noqa: BLE001
            print(f"  ✗ {label}: {e}")
            import traceback; traceback.print_exc()
            failures.append((label, str(e)))
    print()
    if failures:
        print(f"FAILED ({len(failures)}/{len(CHECKS)})")
        return 1
    if args.bless:
        maybe_bless()
    print(f"OK ({len(CHECKS)} checks)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
