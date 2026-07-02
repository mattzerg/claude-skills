#!/usr/bin/env python3
"""viz-review — audit charts + decks + chart specs against named style rules."""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path.home() / ".claude/skills/_consultant/python"))
from _bootstrap_helper import ensure_venv  # noqa: E402

ensure_venv(__file__)

import argparse  # noqa: E402
import json  # noqa: E402
import re  # noqa: E402


# --- chart mode -------------------------------------------------------------

def chart(args) -> int:
    path = Path(args.path)
    if not path.exists():
        print(f"ERROR: chart not found at {path}")
        return 1
    findings = []

    # File size sanity
    size = path.stat().st_size
    if size < 5_000:
        findings.append(("MED", f"File only {size} bytes — likely a rendering failure"))
    elif size > 2_000_000:
        findings.append(("LOW", f"File {size//1024}KB — unusually large for a branded chart"))

    # Caption sidecar
    cap_path = path.with_suffix(".caption.md")
    if not cap_path.exists():
        findings.append(("MED", "No `.caption.md` sidecar — action-title caption missing"))
    else:
        cap_text = cap_path.read_text()
        if "Source:" not in cap_text and "[source:" not in cap_text and not args.no_source_check:
            findings.append(("MED", "Caption sidecar has no `Source:` line — every chart needs sourcing"))
        if cap_text.startswith("# [draft"):
            findings.append(("MED", "Caption is still a `[draft action title]` placeholder"))

    # Suffix check (PNG + SVG paired)
    svg_path = path.with_suffix(".svg")
    if path.suffix == ".png" and not svg_path.exists():
        findings.append(("LOW", "No paired `.svg` vector — chart-builder writes both by default"))

    # Pixel-level Tufte-ish checks via PIL
    try:
        from PIL import Image  # type: ignore
        img = Image.open(path)
        w, h = img.size
        if w < 800 or h < 400:
            findings.append(("LOW", f"Image resolution {w}×{h} — chart-builder default is ≥1200×675"))
        # rough check: aspect ratio
        if h > w * 1.5:
            findings.append(("LOW", f"Aspect ratio {w}:{h} is taller than wide — most consultant charts are landscape"))
    except Exception as e:  # noqa: BLE001
        findings.append(("LOW", f"PIL inspection failed: {e}"))

    return _emit("chart", path, findings)


# --- deck mode --------------------------------------------------------------

KEY_TYPES_BY_NAME = {"title", "exec-summary", "section-divider", "chart", "two-col",
                     "stats-strip", "quote", "table", "recommendation",
                     "appendix-sources", "appendix", "support", "key-section"}


def deck(args) -> int:
    path = Path(args.path)
    if not path.exists():
        print(f"ERROR: deck not found at {path}")
        return 1
    findings = []

    if path.suffix in (".md", ".markdown"):
        slides = _slides_from_storyline(path)
    elif path.suffix == ".pptx":
        slides = _slides_from_pptx(path)
    else:
        print(f"ERROR: unknown extension {path.suffix}; pass a .pptx or storyline.md")
        return 1

    if not slides:
        findings.append(("HIGH", "No slides parsed from the deck/storyline"))
        return _emit("deck", path, findings)

    # Slide count
    n = len(slides)
    if n > 35:
        findings.append(("HIGH", f"{n} slides — over the 35-slide hard cap; trim the storyline"))
    elif n > 25:
        findings.append(("MED", f"{n} slides — over the 25-slide soft cap; consider merging supports"))

    # Required slides
    types = [s.get("type", "support") for s in slides]
    if "title" not in types:
        findings.append(("HIGH", "No `title` slide"))
    if "exec-summary" not in types:
        findings.append(("HIGH", "No `exec-summary` slide"))
    if "recommendation" not in types:
        findings.append(("HIGH", "No `recommendation` slide — every consultant engagement needs the moment-of-truth"))
    if "appendix-sources" not in types and "appendix" not in types:
        findings.append(("MED", "No `appendix-sources` slide — quantitative claims should be sourced"))

    # Title compliance
    for s in slides:
        t = s.get("type", "support")
        if t in ("title", "section-divider"):
            continue
        title = (s.get("title") or "").strip()
        words = len(title.split())
        if words < 5:
            findings.append(("MED", f"Slide {s.get('n')} `{t}`: title {words} words — likely a topic, not an action title: {title[:50]!r}"))
        # Has a verb? simple heuristic: any of common verbs or any -s -ed -ing
        has_verb = bool(re.search(r"\b(is|are|drives?|shows?|grew|fell|leaves|leads|cuts?|adds?|raises?|reduces?|requires?|enables?|blocks?|delivers?|holds?)\b",
                                  title.lower())) or bool(re.search(r"\b\w+(?:ed|ing|s)\b", title.lower()))
        if not has_verb and t not in ("title", "section-divider"):
            findings.append(("LOW", f"Slide {s.get('n')}: title has no obvious verb — action titles need predicates: {title[:50]!r}"))

    # Layout variety
    if n > 6:
        from collections import Counter
        type_counts = Counter(types)
        most_common, mc_count = type_counts.most_common(1)[0]
        if mc_count > 15 and mc_count / n > 0.7:
            findings.append(("MED", f"{mc_count}/{n} slides are `{most_common}` — layout variety low"))

    # Empty support slides
    for s in slides:
        if s.get("type") == "support":
            if not s.get("chart_path") and not s.get("table_md"):
                findings.append(("MED", f"Slide {s.get('n')}: `support` with no chart_path or table_md (body will be empty)"))

    # Source for client mode
    mode = (slides[0].get("mode") if slides else None) or "ops"

    return _emit("deck", path, findings)


def _slides_from_storyline(path: Path) -> list[dict]:
    from consultant_kit import frontmatter  # type: ignore
    fm, _ = frontmatter.parse(path)
    slides = fm.get("slides") or fm.get("extra", {}).get("slides") or []
    return slides


def _slides_from_pptx(path: Path) -> list[dict]:
    """Inspect a PPTX and reconstruct slide-spec hints. Best-effort."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return []
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        spec = {"n": i, "type": "support"}
        # Heuristic: find biggest text on the slide as title; eyebrow if Space Mono
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        texts.append({"text": run.text, "size": (run.font.size.pt if run.font.size else 0),
                                      "bold": run.font.bold, "font": run.font.name})
        if texts:
            biggest = max(texts, key=lambda t: t["size"])
            spec["title"] = biggest["text"]
        # Image presence implies "chart" type
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                spec["type"] = "chart"
                spec["chart_path"] = "embedded"
                break
        out.append(spec)
    return out


# --- recipe mode ------------------------------------------------------------

def recipe(args) -> int:
    from consultant_kit import chart as _chart  # type: ignore
    path = Path(args.spec)
    if not path.exists():
        print(f"ERROR: spec not found at {path}")
        return 1
    spec = json.loads(path.read_text())
    findings = _chart.validate_spec(spec)
    return _emit("recipe", path, findings)


# --- emit -------------------------------------------------------------------

def _emit(mode: str, path: Path, findings: list[tuple[str, str]]) -> int:
    severity_order = {"HIGH": 0, "MED": 1, "LOW": 2}
    findings.sort(key=lambda f: severity_order[f[0]])
    print(f"## viz-review {mode} — {path.name}")
    print()
    if not findings:
        print("✅ no findings — passes Tufte + Zerg deck rules")
        return 0
    for sev, msg in findings:
        print(f"- **{sev}** — {msg}")
    high = sum(1 for s, _ in findings if s == "HIGH")
    return 1 if high else 0


def main() -> int:
    p = argparse.ArgumentParser(prog="viz-review")
    sub = p.add_subparsers(dest="cmd", required=True)

    c = sub.add_parser("chart")
    c.add_argument("path")
    c.add_argument("--mode", choices=("default", "dark"), default="default")
    c.add_argument("--no-source-check", action="store_true")
    c.set_defaults(func=chart)

    d = sub.add_parser("deck")
    d.add_argument("path")
    d.set_defaults(func=deck)

    r = sub.add_parser("recipe")
    r.add_argument("spec")
    r.set_defaults(func=recipe)

    args = p.parse_args()
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
